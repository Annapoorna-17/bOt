import os
import io
import base64
import csv
from typing import List, Dict
from pypdf import PdfReader
import fitz  # PyMuPDF
from PIL import Image
from openai import OpenAI
from pinecone import Pinecone
import re, hashlib

# Import libraries for different document types
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import chardet
except ImportError:
    chardet = None

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
EMBED_MODEL = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-large")
CHAT_MODEL = os.getenv("OPENAI_CHAT_MODEL", "gpt-4o")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
# Support either var name; prefer your PINECONE_INDEX_NAME
PINECONE_INDEX = os.getenv("PINECONE_INDEX_NAME", "bot-multi") or os.getenv("PINECONE_INDEX")

if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY missing")
if not PINECONE_API_KEY or not PINECONE_INDEX:
    raise RuntimeError("Pinecone API key or index name missing")

pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(PINECONE_INDEX)
oai = OpenAI(api_key=OPENAI_API_KEY)

def _read_pdf_text(path: str) -> str:
    reader = PdfReader(path)
    return "\n".join([page.extract_text() or "" for page in reader.pages])

def _pil_image_to_base64(image: Image.Image) -> str:
    """Convert PIL Image to base64 string for API transmission."""
    buffered = io.BytesIO()
    # Convert to RGB if necessary (handles RGBA, CMYK, etc.)
    if image.mode not in ('RGB', 'L'):  # L is grayscale
        image = image.convert('RGB')
    image.save(buffered, format="PNG", optimize=True)
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

def _describe_image(img_base64: str, page_num: int) -> str:
    """
    Use GPT-4 Vision to describe the content of an image.
    Returns a text description that can be embedded alongside PDF text.
    """
    try:
        print(f"DEBUG: Calling GPT-4o Vision for image on page {page_num}")
        response = oai.chat.completions.create(
            model="gpt-4o",  # gpt-4o has vision capabilities
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Describe this image in detail. Include any text, charts, diagrams, "
                                "tables, or visual information. If it contains data or specific details, "
                                "extract them precisely. This will be used for document search."
                            )
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{img_base64}",
                                "detail": "high"
                            }
                        }
                    ]
                }
            ],
            max_tokens=500,
            temperature=0.1
        )
        description = response.choices[0].message.content.strip()
        print(f"DEBUG: Vision API returned description ({len(description)} chars): {description[:100]}...")
        return description
    except Exception as e:
        print(f"Warning: Failed to describe image on page {page_num}: {e}")
        import traceback
        traceback.print_exc()
        return ""

def _extract_images_from_pdf(path: str) -> List[tuple[str, int]]:
    """
    Extract all embedded images from a PDF using PyMuPDF.
    Returns list of tuples: (base64_image, page_number)
    """
    images = []
    try:
        pdf_document = fitz.open(path)
        print(f"DEBUG: Opening PDF with {len(pdf_document)} pages")

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            image_list = page.get_images(full=True)
            print(f"DEBUG: Page {page_num + 1} has {len(image_list)} images")

            # Extract each image from the page
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]

                    # Load image using PIL
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    print(f"DEBUG: Image {img_index} on page {page_num + 1}: {pil_image.width}x{pil_image.height}px, mode={pil_image.mode}")

                    # Skip very small images (likely decorative/icons)
                    if pil_image.width < 50 or pil_image.height < 50:
                        print(f"DEBUG: Skipping small image {img_index} ({pil_image.width}x{pil_image.height}px)")
                        continue

                    # Resize if too large to save API costs
                    max_size = 2048
                    if max(pil_image.width, pil_image.height) > max_size:
                        ratio = max_size / max(pil_image.width, pil_image.height)
                        new_size = (int(pil_image.width * ratio), int(pil_image.height * ratio))
                        pil_image = pil_image.resize(new_size, Image.Resampling.LANCZOS)
                        print(f"DEBUG: Resized image to {new_size[0]}x{new_size[1]}px")

                    img_base64 = _pil_image_to_base64(pil_image)
                    images.append((img_base64, page_num + 1))  # 1-indexed page numbers
                    print(f"DEBUG: Successfully extracted and encoded image {img_index} from page {page_num + 1}")

                except Exception as e:
                    print(f"Warning: Could not extract image {img_index} from page {page_num + 1}: {e}")
                    import traceback
                    traceback.print_exc()
                    continue

        pdf_document.close()
        print(f"DEBUG: Total images extracted: {len(images)}")
        return images
    except Exception as e:
        print(f"Warning: Failed to extract images from PDF: {e}")
        import traceback
        traceback.print_exc()
        return []

def _read_pdf_content_with_images(path: str) -> str:
    """
    Extract both text and visual content from a PDF.
    Returns combined text with image descriptions interspersed.
    """
    print(f"DEBUG: Starting PDF content extraction from {path}")

    # Extract text first
    text_content = _read_pdf_text(path)
    print(f"DEBUG: Extracted {len(text_content)} chars of text")

    # Extract and describe images
    images = _extract_images_from_pdf(path)

    if not images:
        print("DEBUG: No images found in PDF, returning text only")
        return text_content

    print(f"DEBUG: Processing {len(images)} images with Vision API...")
    # Describe each image
    image_descriptions = []
    for img_base64, page_num in images:
        desc = _describe_image(img_base64, page_num)
        if desc:
            image_descriptions.append(f"\n[PAGE {page_num} IMAGE]: {desc}\n")
            print(f"DEBUG: Added description for page {page_num}")

    # Combine text and image descriptions
    all_content = text_content
    if image_descriptions:
        all_content += "\n\n=== VISUAL CONTENT FROM DOCUMENT ===\n"
        all_content += "\n".join(image_descriptions)
        print(f"DEBUG: Final content length: {len(all_content)} chars (text: {len(text_content)}, images: {len(all_content) - len(text_content)})")
    else:
        print("DEBUG: No image descriptions generated")

    return all_content

def _read_docx_text(path: str) -> str:
    """Extract text from DOCX files."""
    if DocxDocument is None:
        raise RuntimeError("python-docx library not installed. Install with: pip install python-docx")

    print(f"DEBUG: Extracting text from DOCX: {path}")
    doc = DocxDocument(path)

    # Extract text from paragraphs
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

    # Extract text from tables
    tables_text = []
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells])
            if row_text.strip():
                tables_text.append(row_text)

    content = "\n".join(paragraphs)
    if tables_text:
        content += "\n\n=== TABLES ===\n" + "\n".join(tables_text)

    print(f"DEBUG: Extracted {len(content)} chars from DOCX")
    return content

def _read_xlsx_text(path: str) -> str:
    """Extract text from XLSX files."""
    if load_workbook is None:
        raise RuntimeError("openpyxl library not installed. Install with: pip install openpyxl")

    print(f"DEBUG: Extracting text from XLSX: {path}")
    wb = load_workbook(path, data_only=True)

    all_content = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        all_content.append(f"\n=== SHEET: {sheet_name} ===\n")

        for row in sheet.iter_rows(values_only=True):
            # Filter out None and empty values, convert to string
            row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
            if row_values:
                all_content.append(" | ".join(row_values))

    content = "\n".join(all_content)
    print(f"DEBUG: Extracted {len(content)} chars from XLSX")
    return content

def _read_pptx_text(path: str) -> str:
    """Extract text from PPTX files."""
    if Presentation is None:
        raise RuntimeError("python-pptx library not installed. Install with: pip install python-pptx")

    print(f"DEBUG: Extracting text from PPTX: {path}")
    prs = Presentation(path)

    all_content = []
    for slide_num, slide in enumerate(prs.slides, 1):
        all_content.append(f"\n=== SLIDE {slide_num} ===\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_content.append(shape.text)

            # Extract text from tables in slides
            if shape.shape_type == 19:  # Table
                try:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            all_content.append(row_text)
                except:
                    pass

    content = "\n".join(all_content)
    print(f"DEBUG: Extracted {len(content)} chars from PPTX")
    return content

def _read_txt_text(path: str) -> str:
    """Extract text from plain text files with encoding detection."""
    print(f"DEBUG: Extracting text from TXT: {path}")

    # Try to detect encoding
    encoding = 'utf-8'
    if chardet:
        with open(path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'
            print(f"DEBUG: Detected encoding: {encoding}")

    try:
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()
    except UnicodeDecodeError:
        # Fallback to utf-8 with error handling
        print(f"DEBUG: Failed with {encoding}, falling back to utf-8 with error handling")
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

    print(f"DEBUG: Extracted {len(content)} chars from TXT")
    return content

def _read_csv_text(path: str) -> str:
    """Extract text from CSV files."""
    print(f"DEBUG: Extracting text from CSV: {path}")

    # Try to detect encoding
    encoding = 'utf-8'
    if chardet:
        with open(path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'
            print(f"DEBUG: Detected encoding: {encoding}")

    all_content = []
    try:
        with open(path, 'r', encoding=encoding, newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join([cell.strip() for cell in row if cell.strip()])
                if row_text:
                    all_content.append(row_text)
    except UnicodeDecodeError:
        # Fallback to utf-8 with error handling
        print(f"DEBUG: Failed with {encoding}, falling back to utf-8 with error handling")
        with open(path, 'r', encoding='utf-8', errors='ignore', newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                row_text = " | ".join([cell.strip() for cell in row if cell.strip()])
                if row_text:
                    all_content.append(row_text)

    content = "\n".join(all_content)
    print(f"DEBUG: Extracted {len(content)} chars from CSV")
    return content

def _extract_document_content(path: str, file_extension: str) -> str:
    """
    Extract text content from various document formats.
    Returns the extracted text content.
    """
    ext = file_extension.lower()

    print(f"DEBUG: Extracting content from {ext} file: {path}")

    if ext == '.pdf':
        return _read_pdf_content_with_images(path)
    elif ext in ['.docx', '.doc']:
        return _read_docx_text(path)
    elif ext in ['.xlsx', '.xls']:
        return _read_xlsx_text(path)
    elif ext in ['.pptx', '.ppt']:
        return _read_pptx_text(path)
    elif ext == '.csv':
        return _read_csv_text(path)
    elif ext in ['.txt', '.md', '.rst', '.log']:
        return _read_txt_text(path)
    else:
        # For unknown types, try reading as text
        print(f"DEBUG: Unknown extension {ext}, attempting text extraction")
        try:
            return _read_txt_text(path)
        except Exception as e:
            raise RuntimeError(f"Unsupported file type: {ext}. Error: {e}")

def _chunk_text(text: str, max_chars: int = 3000, overlap: int = 400):
    """
    Simple, fast chunker that respects sentence boundaries when possible.
    - max_chars ~ roughly 750 tokens (4 chars/token heuristic), adjust as needed.
    - overlap keeps a bit of context between chunks.
    """
    text = re.sub(r'\s+', ' ', text).strip()

    if len(text) <= max_chars:
        return [text] if text else []

    chunks = []
    start = 0
    end = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        # try to cut at the last sentence boundary within the window
        window = text[start:end]
        cut = max(window.rfind('. '), window.rfind('? '), window.rfind('! '))
        if cut == -1 or end == len(text):
            cut = len(window)
        chunk = window[:cut].strip()
        if chunk:
            chunks.append(chunk)
        # move start forward with overlap
        start = start + cut
        start = max(0, start - overlap)

        # avoid infinite loop on very small residuals
        if end == len(text):
            break

    # dedupe / clean
    out = []
    last = ""
    for c in chunks:
        if c and c != last:
            out.append(c)
            last = c
    return out

def embed_chunks(chunks: List[str]) -> List[List[float]]:
    resp = oai.embeddings.create(model=EMBED_MODEL, input=chunks)
    return [d.embedding for d in resp.data]

def upsert_chunks(tenant_code: str, user_code: str, doc_filename: str, chunks: List[str]) -> int:
    embs = embed_chunks(chunks)
    vectors = []
    for i, (chunk, vec) in enumerate(zip(chunks, embs)):
        vid = hashlib.sha256(f"{tenant_code}:{doc_filename}:{i}".encode()).hexdigest()
        vectors.append({
            "id": vid,
            "values": vec,
            "metadata": {
                "tenant_code": tenant_code,
                "user_code": user_code,
                "source_type": "document",  # Distinguish from websites
                "doc": doc_filename,
                "chunk_index": i,
                "text": chunk
            }
        })
    index.upsert(vectors=vectors, namespace=tenant_code)
    return len(vectors)

def document_to_pinecone(file_path: str, tenant_code: str, user_code: str, stored_filename: str) -> int:
    """
    Extract content from any supported document format and index to Pinecone.
    Supports: PDF, DOCX, XLSX, PPTX, CSV, TXT, MD, and more.
    """
    # Get file extension from stored filename
    _, file_extension = os.path.splitext(stored_filename)

    # Extract content based on file type
    content = _extract_document_content(file_path, file_extension)

    if not content.strip():
        print(f"WARNING: No content extracted from {stored_filename}")
        return 0

    chunks = _chunk_text(content)
    return upsert_chunks(tenant_code, user_code, stored_filename, chunks)

# Keep backward compatibility
def pdf_to_pinecone(file_path: str, tenant_code: str, user_code: str, stored_filename: str) -> int:
    """Backward compatibility wrapper for document_to_pinecone."""
    return document_to_pinecone(file_path, tenant_code, user_code, stored_filename)

def search(tenant_code: str, query: str, top_k: int = 8, filter_user_code: str | None = None, source_type: str = "all"):
    """
    Search for relevant content in Pinecone.

    Args:
        tenant_code: Tenant identifier
        query: Search query
        top_k: Number of results to return
        filter_user_code: Optional user filter
        source_type: Filter by source type - "all", "documents", or "websites"
    """
    q_emb = oai.embeddings.create(model=EMBED_MODEL, input=[query]).data[0].embedding

    # Build filter - start with tenant filter
    flt = {"tenant_code": {"$eq": tenant_code}}

    # Add source type filter if not "all"
    if source_type == "documents":
        flt = {"$and": [flt, {"source_type": {"$eq": "document"}}]}
        print(f"DEBUG: Filtering to DOCUMENTS only")
    elif source_type == "websites":
        flt = {"$and": [flt, {"source_type": {"$eq": "website"}}]}
        print(f"DEBUG: Filtering to WEBSITES only")
    else:
        print(f"DEBUG: Searching ALL sources (documents + websites)")

    # Add user filter if specified
    if filter_user_code:
        flt = {"$and": [flt, {"user_code": {"$eq": filter_user_code}}]}

    print(f"DEBUG: Pinecone filter: {flt}")

    res = index.query(
        vector=q_emb, top_k=top_k, namespace=tenant_code,
        filter=flt, include_metadata=True
    )

    matches = res.matches or []
    print(f"DEBUG: Found {len(matches)} matches")

    # Debug: Show what sources were found
    if matches:
        for i, m in enumerate(matches[:5]):  # Show first 5
            source_type_found = m.metadata.get("source_type", "unknown")
            if source_type_found == "website":
                source_name = m.metadata.get("url", "unknown")
            else:
                source_name = m.metadata.get("doc", "unknown")
            score = m.score if hasattr(m, 'score') else 'N/A'
            print(f"DEBUG: Match {i+1}: {source_type_found} - {source_name} (score: {score})")

    return matches

def _clean_answer_text(text: str) -> str:
    """
    Clean up answer text by removing unwanted symbols and formatting artifacts.
    Enhanced for chatbot-friendly responses with better list handling.
    """
    # Remove excessive newlines (keep max 2 consecutive newlines)
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Remove weird unicode characters and symbols that might appear
    text = re.sub(r'[\u200b-\u200d\ufeff]', '', text)  # Zero-width spaces

    # Remove excessive whitespace
    text = re.sub(r'[ \t]+', ' ', text)

    # Remove markdown formatting symbols (bold, italic, code)
    # Remove ** for bold
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    # Remove __ for bold
    text = re.sub(r'__(.+?)__', r'\1', text)
    # Remove * for italic
    text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'\1', text)
    # Remove _ for italic (but not in words)
    text = re.sub(r'(?<!\w)_(?!_)(.+?)(?<!_)_(?!\w)', r'\1', text)
    # Remove backticks for code
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove ~~strikethrough~~
    text = re.sub(r'~~(.+?)~~', r'\1', text)

    # Remove markdown headers (##, ###, etc.)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)

    # Clean up bullet points and list markers - standardize to •
    text = re.sub(r'^\s*[-*+]\s+', '• ', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*[•●○◦]\s+', '• ', text, flags=re.MULTILINE)

    # Clean up numbered lists - ensure proper format (number followed by period and space)
    text = re.sub(r'^\s*(\d+)[.)]\s+', r'\1. ', text, flags=re.MULTILINE)

    # Remove excessive spaces before punctuation
    text = re.sub(r'\s+([.,;:!?])', r'\1', text)

    # Remove brackets around citations like [1], [2], etc.
    text = re.sub(r'\[\d+\]', '', text)

    # Remove HTML-like tags if present
    text = re.sub(r'<[^>]+>', '', text)

    # Clean up phrases that sound too formal for a chat
    text = re.sub(r'^Based on the (provided )?context(,| provided)?\s*', '', text, flags=re.IGNORECASE)
    text = re.sub(r'^According to the (provided )?context(,| provided)?\s*', '', text, flags=re.IGNORECASE)
    text = re.sub(r'^From the context(,| provided)?\s*', '', text, flags=re.IGNORECASE)

    # Remove any trailing/leading whitespace
    text = text.strip()

    return text

def synthesize_answer(question: str, contexts: List[str]) -> str:
    """
    Generate a refined answer using the provided contexts.
    The answer is cleaned to remove unwanted symbols and formatting artifacts.
    Enhanced with chatbot-friendly formatting instructions.
    """
    # Enhanced system prompt for chatbot-friendly responses
    system_prompt = """You are a helpful AI assistant in a chatbot interface. Provide clear, accurate answers.

FORMATTING RULES:
1. Use natural, conversational language suitable for chat
2. Never use markdown symbols like **, __, ~~, or ` unless absolutely necessary
3. Avoid unnecessary special characters, emojis, or decorative symbols
4. Keep responses clean and easy to read in a chat window
5. Use proper sentence structure with correct punctuation

WHEN LISTS ARE EXPECTED:
- If the question asks for a list (e.g., "list all", "what are the", "give me the"), format as a clear numbered or bulleted list
- Use simple bullet points (•) or numbers (1., 2., 3.)
- Each item should be on a new line
- Keep list items concise and scannable
- Example format:
  • Item one
  • Item two
  • Item three

RESPONSE STYLE:
- Be direct and precise
- Avoid verbose introductions like "Based on the context provided..."
- For factual questions, give factual answers directly
- For list questions, give list answers
- Answer based on the context provided
- If information is clearly missing from the context, politely say you don't have that information

Remember: You are chatting with a user, not writing a formal document."""

    user_prompt = (
        f"Question: {question}\n\n"
        "Context:\n" + "\n\n---\n".join(contexts[:12]) + "\n\n"
        "Provide a clear, helpful answer based on the information in the context above."
    )

    chat = oai.chat.completions.create(
        model=CHAT_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
    )
    raw_answer = chat.choices[0].message.content.strip()

    # Clean up the answer
    cleaned_answer = _clean_answer_text(raw_answer)

    return cleaned_answer
